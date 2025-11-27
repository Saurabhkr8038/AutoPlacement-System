from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "University Placement Automation System"
    subtitle.text = "Streamlining Recruitment. Connecting Talent.\n\nTeam Members:\nSaurabh Kumar (saurabh.24scse1180099@galgotiasuniversity.ac.in)\nYashdeep Shukla (Yashdeep.24scse1180255@galgotiasuniversity.ac.in)"

    # Helper function to add a bullet slide
    def add_bullet_slide(title_text, content_items):
        slide_layout = prs.slide_layouts[1] # Title and Content
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        tf = slide.placeholders[1].text_frame
        for item in content_items:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0

    # Slide 2: Problem Statement
    add_bullet_slide("Problem Statement", [
        "Manual Data Handling: Reliance on error-prone Excel sheets.",
        "Inefficient Filtering: Time-consuming manual sorting of students by CGPA.",
        "Data Redundancy: Lack of centralized tracking for eligibility.",
        "Communication Gap: No unified platform for job postings."
    ])

    # Slide 3: Proposed Solution
    add_bullet_slide("Proposed Solution", [
        "Centralized Java Application: Secure desktop system for data management.",
        "Automated Filtering Engine: Instant shortlisting algorithms.",
        "Database Integration: MySQL for persistent storage.",
        "Role-Based Access: Separate modules for Admins and Students."
    ])

    # Slide 4: Key Features
    add_bullet_slide("Key Features", [
        "Student Management: CRUD operations for profiles.",
        "Smart Filtering: Instant candidate list generation.",
        "Java Backend: Core Java and OOP principles.",
        "JDBC Connectivity: Live MySQL synchronization.",
        "Multithreading: Background data loading.",
        "Data Security: Encapsulated data models."
    ])

    # Slide 5: System Architecture
    add_bullet_slide("System Architecture", [
        "Presentation Layer: Console/GUI Interface.",
        "Service Layer: Filtering logic (PlacementService).",
        "DAO Layer: Database operations (StudentDAO, CompanyDAO).",
        "Database Layer: MySQL Database.",
        "Utilities: DBConnection (Singleton), Thread Managers."
    ])

    # Helper function for code slides
    def add_code_slide(title_text, description, code_snippet):
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        # Add description
        tf = slide.placeholders[1].text_frame
        p = tf.add_paragraph()
        p.text = description
        p.font.size = Pt(18)
        
        # Add code box
        left = Inches(0.5)
        top = Inches(2.5)
        width = Inches(9)
        height = Inches(4.5)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        tf_code = textbox.text_frame
        tf_code.text = code_snippet
        tf_code.paragraphs[0].font.name = "Courier New"
        tf_code.paragraphs[0].font.size = Pt(12)

    # Slide 6: OOP Concepts Used
    oop_code = """public class Student extends User {
    private double cgpa; // Encapsulation

    public Student(...) { super(...); } // Inheritance

    @Override
    public String toString() { ... } // Polymorphism
}"""
    add_code_slide("OOP Concepts Used", 
                   "Encapsulation, Inheritance, Polymorphism, Interfaces.",
                   oop_code)

    # Slide 7: Collections & Generics
    collections_code = """public class PlacementService {
    private Map<String, List<Student>> companyShortlists;

    public PlacementService() {
        this.companyShortlists = new HashMap<>();
    }
    // ...
}"""
    add_code_slide("Collections & Generics", 
                   "Used ArrayList for dynamic lists and HashMap for lookups.",
                   collections_code)

    # Slide 8: Multithreading
    thread_code = """public class LoaderThread extends Thread {
    @Override
    public void run() {
        System.out.println("Background loading...");
        // Fetch heavy data
    }
}
// Main:
LoaderThread loader = new LoaderThread();
loader.start();"""
    add_code_slide("Multithreading", 
                   "Background data loading to keep UI responsive.",
                   thread_code)

    # Slide 9: Database Layer (JDBC)
    jdbc_code = """String sql = "INSERT INTO students ... VALUES (?, ?)";
try (PreparedStatement pstmt = conn.prepareStatement(sql)) {
    pstmt.setString(1, name);
    pstmt.executeUpdate();
} catch (SQLException e) {
    e.printStackTrace();
}"""
    add_code_slide("Database Layer (JDBC)", 
                   "Secure database connectivity using PreparedStatement.",
                   jdbc_code)

    # Slide 10: Workflow Diagram
    add_bullet_slide("Workflow", [
        "1. Login: Admin authenticates.",
        "2. Input: Admin sets criteria (e.g., CGPA > 7.5).",
        "3. Process: System queries DB & filters students.",
        "4. Output: Shortlisted candidates displayed.",
        "5. Storage: Updates saved to MySQL."
    ])

    # Slide 11: Conclusion & Future Scope
    add_bullet_slide("Conclusion & Future Scope", [
        "Conclusion: Successfully automates placement workflow using Java OOP, JDBC, and Multithreading.",
        "Future Scope:",
        "- Web Portal (JSP/Servlet)",
        "- AI Resume Parsing",
        "- Email Automation"
    ])

    prs.save('University_Placement_Automation_System.pptx')
    print("Presentation saved successfully.")

if __name__ == "__main__":
    create_presentation()
